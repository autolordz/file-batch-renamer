# Copyright (c) 2018 Autoz https://github.com/autolordz

# Permission is hereby granted, free of charge, to any person obtaining a copy
# of this software and associated documentation files (the "Software"), to deal
# in the Software without restriction, including without limitation the rights
# to use, copy, modify, merge, publish, distribute, sublicense, and/or sell
# copies of the Software, and to permit persons to whom the Software is
# furnished to do so, subject to the following conditions:

# The above copyright notice and this permission notice shall be included in all
# copies or substantial portions of the Software.

# THE SOFTWARE IS PROVIDED "AS IS", WITHOUT WARRANTY OF ANY KIND, EXPRESS OR
# IMPLIED, INCLUDING BUT NOT LIMITED TO THE WARRANTIES OF MERCHANTABILITY,
# FITNESS FOR A PARTICULAR PURPOSE AND NONINFRINGEMENT. IN NO EVENT SHALL THE
# AUTHORS OR COPYRIGHT HOLDERS BE LIABLE FOR ANY CLAIM, DAMAGES OR OTHER
# LIABILITY, WHETHER IN AN ACTION OF CONTRACT, TORT OR OTHERWISE, ARISING FROM,
# OUT OF OR IN CONNECTION WITH THE SOFTWARE OR THE USE OR OTHER DEALINGS IN THE
# SOFTWARE.

# -*- coding: utf-8 -*-
"""
Created on Mon Nov 26 11:56:28 2018

@author: zhz
"""

#%%
import zhon.hanzi,zhon.cedict
import os,re,io,glob,shutil,string,platform
import itertools as it
from pdfminer.high_level import extract_text_to_fp
import pytesseract
from PIL import Image
from docx import Document
from pptx import Presentation
from xlrd import open_workbook
from win32com.client import Dispatch # for office 97-2003

#%%
def parse_subpath(path,file):
    '''make subpath'''
    if not path: return file
    if not os.path.exists(path):
        os.mkdir(path)
    return os.path.join(path,file)


def os_rename(origin,dist):
    if dist:
        header = os.path.dirname(origin)+'\\' if os.path.dirname(origin) else ''
        file_n = header + dist + os.path.splitext(origin)[1]
        if not file_n == origin:
            try:
                os.rename(origin,file_n)
            except FileExistsError:
                os_rename(origin,dist+'_copy')
            print('>>> 重命名: 【%s】=>【%s】'%(origin,file_n))

def clean_txt_func(x,**kwargs):
    x = re.sub(r'\s+',',',x)
    # x = unicodedata.normalize('NFKC', x)
    s = int(kwargs.get('txt_l',len(''.join(x))))
    punc_all = string.punctuation + zhon.hanzi.punctuation
    char_all = string.printable + zhon.cedict.all
    if re.search(r'[\u4e00-\u9fff]+',x):
        x = re.sub(r'(?<=[%s\w]{2})[%s]'%(zhon.cedict.all,punc_all),'|',x)
        x = re.sub(r'[%s]'%punc_all.replace('|',''),'',x) # chinese punctuation
        x = re.sub(r'[^%s]'%char_all,'',x)
    else:
        x = re.sub(r'(?<=\w{2})[%s]'%(string.punctuation),'|',x)
        x = re.sub(r'[%s]'%string.punctuation.replace('|',''),'',x) # punctuation
        x = re.sub(r'[^%s]'%string.printable,'',x) # printable
    x = re.sub(r'PowerPoint|演示文稿|Sheet1','',x) # PowerPoint Excel
    # re.compile(r'^((?!第.*卷|期|编|cid).)*$').search,x)
    xx = re.split(r'\|',x)
    xx = '_'.join(xx)[:s]
    return xx


#%% rename office,officex

def rename_officex(file,**kwargs):
    '''rename only judgment doc files'''
    suffix = os.path.splitext(file)[1]

    if suffix == '.docx':
        try:
            doc = Document(file)
            x = ','.join([para.text.strip() for para in doc.paragraphs])
            x = clean_txt_func(x,**kwargs)
            print('>>> 找到 %s 内容: %s'%(file,x))
            os_rename(file,x)
            return x
        except Exception as e:
            print('>>> 读取 %s 失败,可能格式不正确 => %s'%(file,e))

    if suffix == '.pptx':
        try:
            prs = Presentation(file)
            txt = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            txt.append(run.text)
            x = clean_txt_func(','.join(txt),**kwargs)
            print('>>> 找到 %s 内容: %s'%(file,x))
            os_rename(file,x)
            return x
        except Exception as e:
            print('>>> 读取 %s 失败,可能格式不正确 => %s'%(file,e))

    if suffix in ['.xlsx','.xls']:
        try:
            exl = open_workbook(file)
            txt = []
            for sheet in exl.sheets():
                for i in range(sheet.nrows):
                    for j in range(sheet.ncols):
                        v = sheet.cell_value(i,j)
                        if not v:continue
                        txt.append(v)
            x = clean_txt_func(','.join(txt),**kwargs)
            print('>>> 找到 %s 内容: %s'%(file,x))
            os_rename(file,x)
            return x
        except Exception as e:
            print('>>> 读取 %s 失败,可能格式不正确 => %s'%(file,e))
    return None

def get_txt_text(file,**kwargs):
    try:
        with open(file,'r',errors='ignore') as f:
            x = ','.join([line.rstrip('\n') for line in f])
            x = clean_txt_func(x,**kwargs)
            return x
    except Exception as e:
        print('>>> 读取 %s 失败,可能格式不正确 => %s'%(file,e))
    return ''

def rename_office(file,**kwargs):
    name = os.path.splitext(file)[0]
    suffix = os.path.splitext(file)[1]

    if suffix == '.txt':
        x = get_txt_text(file,**kwargs)
        print('>>> 找到 %s 内容: %s'%(file,x))
        os_rename(file,x)

    if suffix == '.doc':
        file_txt = name + '_doc.txt'
        word = Dispatch("Word.Application")
        word.Documents.Open(os.path.abspath(file))
        word.ActiveDocument.SaveAs(os.path.abspath(file_txt),FileFormat=3)
        word.ActiveDocument.Close()
        x = get_txt_text(file_txt,**kwargs)
        print('>>> 找到 %s 内容: %s'%(file,x))
        os_rename(file,x)
        os.remove(file_txt)

    if suffix == '.ppt':
        txt = []
        try:
            app = Dispatch("PowerPoint.Application")
        except Exception as e:
            print(e)
            # print('尝试使用Tika来读取')
            # from tika import parser
            # txt = parser.from_file(file)['content']
        else:
            ppt = app.Presentations.Open(os.path.abspath(file))
            for slide in ppt.Slides:
                for shape in slide.Shapes:
                    if not shape.HasTextFrame:
                        continue
                    txt.append(shape.TextFrame.TextRange.Text)
            app.ActivePresentation.Close()
        x = clean_txt_func(','.join(txt),**kwargs)
        print('>>> 找到 %s 内容: %s'%(file,x))
        os_rename(file,x)

    if suffix == '.xls':
        try:
            app = Dispatch("Excel.Application")
        except Exception as e:
            print(e)
            rename_officex(file,**kwargs)
            return True
        exl = app.Workbooks.Open(os.path.abspath(file))
        txt = []
        for sheet in exl.Sheets:
            sheet = sheet.Range(sheet.Cells(1, 1), sheet.Cells(20, 10))
            for cell in sheet.Cells:
                if not cell.Value:
                    continue
                txt.append(cell.Value)
        app.ActiveWorkbook.Close()
        x = clean_txt_func(','.join(txt),**kwargs)
        print('>>> 找到 %s 内容: %s'%(file,x))
        os_rename(file,x)

    return True

#%% rename image

def get_image_txt(file,**kwargs):
    try:
        img_h = kwargs.get('img_h',1)
        print('image:',file)
        img = Image.open(file)
        if img.width > img.height: img = img.rotate(90,expand=1)
        print('image size :',img.size)
        img = img.crop((0,0,img.width,img.height/img_h))
        print('image size 2:',img.size)

        pytesseract.pytesseract.tesseract_cmd = 'c:\\Program Files (x86)\\Tesseract-OCR\\tesseract.exe' \
        if '64bit' in platform.architecture() else 'c:\\Program Files\\Tesseract-OCR\\tesseract.exe'

        x = pytesseract.image_to_string(img,lang='chi_sim') # eng
        x = re.sub(r'\s+',',',x)
        print('>>> 解析 %s \n 内容: %s'%(file,x))
        return x
    except Exception as e:
        print(e)
    return ''

def rename_image(file,**kwargs):
    x = get_image_txt(file,**kwargs)
    x = clean_txt_func(x,**kwargs)
    os_rename(file,x)
    return True

#%% rename pdf

def get_pdf_txt(ifile,**kwargs):
    txt='';odir = ''.join(['.\\',os.path.splitext(ifile)[0],'_img'])
    with io.StringIO() as ofp:
        with open(ifile, "rb") as ifp:
            try:
                extract_text_to_fp(ifp,ofp,maxpages=2)
                                   # password=''
                                   # output_dir=odir,
                                   # page_numbers=set([0]),
            except Exception as e:
                print('error:',e.__class__.__name__)
        txt = ofp.getvalue()
        print('full len:',len(txt))
        if len(txt) < 10:
            print('====decode images===')
            extrectImage.main(sourceName=ifile,outputFolder=odir,**kwargs)
            subext = [parse_subpath(odir,x) for x in
                      ['*.png','*.jpg','*.jpeg','*.bmp','*.tif']]
            images = list(it.chain(*(glob.iglob(e) for e in subext)))
            print(images)
            if images: #decode the 1st page
                txt = ''
                for image in images:
                    if len(txt) > 10:break
                    txt = get_image_txt(image,**kwargs)
    return txt,odir

def rename_pdf(file,**kwargs):
    del_img_dir = kwargs.get('del_img_dir',True)
    print('pdf name: ', file)
    x,odir = get_pdf_txt(file,**kwargs)
    x = clean_txt_func(x,**kwargs)
    print('cut txt:',x)
    print('cut len:',len(x))
    if del_img_dir: shutil.rmtree(odir)
    os_rename(file,x)
    return True

#%% main

def get_files_list(suffix,subdir=None):
    suffix = [parse_subpath(subdir,x) for x in suffix]
    files = list(it.chain(*(glob.iglob(e) for e in suffix)))
    return files

suffix = ['*.doc','*.ppt','*.xls','*.txt',
         '*.docx','*.pptx','*.xlsx','*.pdf',
         '*.png','*.jpg','*.jpeg','*.bmp','*.tif']

files = get_files_list(suffix)

if len(files)>0:
    print('>>> 正在重命名 ...')
    print('================ \n')
    for i,file in enumerate(files):
        if '~$' in file: continue
        if os.path.splitext(file)[1] in ['.doc','.ppt','.xls','.txt']:
            rename_office(file,txt_l = 30)
        elif os.path.splitext(file)[1] in ['.docx','.pptx','.xlsx']:
            rename_officex(file,txt_l = 30)
        elif os.path.splitext(file)[1] in ['.pdf']:
            rename_pdf(file,txt_l = 30,img_h = 3,num_pages=1,del_img_dir = False)
        elif os.path.splitext(file)[1] in ['.png','.jpg','.jpeg','.bmp','.tif']:
            rename_image(file,txt_l = 30,img_h = 3)
        else:
            print(' 跳过文件 \n')
        print('\n ================ \n')
    print('>>> 重命名完毕...')

